from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from deck_utils import THEMES, add_footer, add_header, add_lines, add_panel, add_text, new_presentation, rgb, save, set_background


def build_template(template_id: str, destination: Path) -> Path:
    theme = THEMES[template_id]
    prs = new_presentation()
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    set_background(cover, theme.background)
    rule = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, 180000)
    rule.fill.solid()
    rule.fill.fore_color.rgb = rgb(theme.accent)
    rule.line.fill.background()
    add_text(cover, theme.name, 0.75, 1.45, 10.8, 0.78, size=34, color=theme.ink, bold=True)
    add_text(cover, "[主题或一句话结论]", 0.78, 2.42, 10.5, 0.6, size=20, color=theme.muted)
    add_text(cover, "[汇报对象]  |  [日期]", 0.8, 5.95, 7.0, 0.35, size=10, color=theme.muted)
    add_text(cover, "MARKET DIRECTOR", 10.0, 5.95, 2.5, 0.35, size=9, color=theme.accent, bold=True, align=PP_ALIGN.RIGHT)

    section = prs.slides.add_slide(blank)
    set_background(section, theme.ink)
    add_text(section, "01", 0.75, 1.15, 1.2, 0.6, size=16, color=theme.accent, bold=True)
    add_text(section, "[章节标题]", 0.75, 2.0, 10.8, 0.85, size=34, color=theme.background, bold=True)
    add_text(section, "[这一章节要回答的问题]", 0.78, 3.0, 10.2, 0.5, size=16, color=theme.pale)

    content = prs.slides.add_slide(blank)
    set_background(content, theme.background)
    add_header(content, "[结论式页面标题]", "[一句话说明本页为什么重要]", theme, 3)
    add_panel(content, 0.68, 1.65, 7.6, 4.95, "FFFFFF", theme.pale)
    add_text(content, "[主要证据或图表区域]", 1.0, 2.0, 7.0, 0.45, size=18, color=theme.ink, bold=True)
    add_lines(content, ["[关键事实 1]", "[关键事实 2]", "[关键事实 3]"], 1.0, 2.75, 6.8, 2.5, size=14, color=theme.ink, bullet_color=theme.accent)
    add_panel(content, 8.6, 1.65, 4.05, 4.95, theme.pale)
    add_text(content, "管理含义", 8.95, 2.0, 3.3, 0.4, size=15, color=theme.ink, bold=True)
    add_lines(content, ["[判断]", "[建议行动]", "[责任人与期限]"], 8.95, 2.7, 3.2, 2.9, size=12, color=theme.ink, bullet_color=theme.secondary)
    add_footer(content, "来源：[文件或链接]", theme)

    ending = prs.slides.add_slide(blank)
    set_background(ending, theme.background)
    add_text(ending, "下一步", 0.75, 1.2, 5.0, 0.7, size=32, color=theme.ink, bold=True)
    add_text(ending, "[需要确认的决策与行动]", 0.78, 2.1, 8.5, 0.5, size=17, color=theme.muted)
    add_lines(ending, ["[行动 1]", "[行动 2]", "[行动 3]"], 0.8, 3.0, 8.4, 2.4, size=16, color=theme.ink, bullet_color=theme.accent)
    add_text(ending, "谢谢", 10.1, 5.9, 2.2, 0.5, size=20, color=theme.accent, bold=True, align=PP_ALIGN.RIGHT)

    prs.core_properties.title = theme.name
    prs.core_properties.subject = "可编辑市场演示文稿模板"
    prs.core_properties.author = "市场总监工作台"
    return save(prs, destination)


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate the editable market template library.")
    parser.add_argument("--project", type=Path, default=Path.cwd())
    parser.add_argument("--plugin-assets", type=Path)
    args = parser.parse_args()

    project = args.project.resolve()
    project_dir = project / "library" / "templates"
    plugin_dir = args.plugin_assets.resolve() if args.plugin_assets else None
    created = []
    for template_id in THEMES:
        target = project_dir / f"{template_id}.pptx"
        build_template(template_id, target)
        created.append(str(target))
        if plugin_dir:
            build_template(template_id, plugin_dir / f"{template_id}.pptx")

    catalog_path = project_dir / "template-catalog.json"
    catalog = json.loads(catalog_path.read_text(encoding="utf-8"))
    catalog["generated_files"] = created
    catalog_path.write_text(json.dumps(catalog, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps({"created": created}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

