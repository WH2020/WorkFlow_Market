from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


def validate_deck(path: Path, min_slides: int | None = None, max_slides: int | None = None) -> list[str]:
    errors: list[str] = []
    prs = Presentation(path)
    count = len(prs.slides)
    if min_slides is not None and count < min_slides:
        errors.append(f"slide_count_below_min:{count}<{min_slides}")
    if max_slides is not None and count > max_slides:
        errors.append(f"slide_count_above_max:{count}>{max_slides}")
    for index, slide in enumerate(prs.slides, 1):
        if not slide.shapes:
            errors.append(f"empty_slide:{index}")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text:
                continue
            if len(text) > 1200:
                errors.append(f"dense_text:{index}:{len(text)}")
            if "[待" in text or "[主题" in text or "[结论" in text:
                errors.append(f"placeholder_text:{index}:{text[:80]}")
    return errors


def main() -> int:
    parser = argparse.ArgumentParser(description="Validate generated PPTX and JSON outputs.")
    parser.add_argument("--pptx", type=Path)
    parser.add_argument("--json", type=Path)
    parser.add_argument("--min-slides", type=int)
    parser.add_argument("--max-slides", type=int)
    args = parser.parse_args()
    errors: list[str] = []
    if args.pptx:
        errors.extend(validate_deck(args.pptx.resolve(), args.min_slides, args.max_slides))
    if args.json:
        try:
            json.loads(args.json.resolve().read_text(encoding="utf-8"))
        except Exception as exc:  # noqa: BLE001
            errors.append(f"invalid_json:{exc}")
    print(json.dumps({"ok": not errors, "errors": errors}, ensure_ascii=False, indent=2))
    return 0 if not errors else 1


if __name__ == "__main__":
    raise SystemExit(main())

