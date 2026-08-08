from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333333)
SLIDE_HEIGHT = Inches(7.5)
FONT_HEAD = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"


@dataclass(frozen=True)
class Theme:
    key: str
    name: str
    background: str
    ink: str
    muted: str
    accent: str
    secondary: str
    warning: str
    pale: str


THEMES = {
    "ceo-weekly": Theme("ceo-weekly", "CEO 市场周报", "F7F7F5", "151719", "667079", "C92A3A", "147D64", "D99025", "E8ECEA"),
    "government-formal": Theme("government-formal", "地方政府合作方案", "FAF9F7", "1A1A1A", "6D625E", "A61B1B", "B38A2E", "B05C2E", "EFE9E5"),
    "customer-solution": Theme("customer-solution", "客户解决方案", "F8FAFB", "172126", "65747C", "147D64", "2563A6", "D97732", "E7EFEE"),
    "industry-research": Theme("industry-research", "行业研究报告", "FAFAF8", "1A2421", "66706C", "167C6A", "D49B24", "D65A4A", "E7EFE9"),
    "sales-review": Theme("sales-review", "销售复盘报告", "F8F9F7", "17201C", "67706B", "1F8066", "D44A4A", "D08A27", "E6ECE8"),
}


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = "151719",
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
    font: str = FONT_BODY,
    line_spacing: float = 1.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_lines(
    slide,
    lines: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str,
    bullet_color: str,
    gap: float = 0.12,
):
    lines = [str(line) for line in lines if str(line).strip()]
    if not lines:
        lines = ["暂无内容"]
    row_h = max(0.38, (h - gap * max(0, len(lines) - 1)) / len(lines))
    for index, line in enumerate(lines):
        top = y + index * (row_h + gap)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(top + 0.13), Inches(0.09), Inches(0.09))
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb(bullet_color)
        circle.line.fill.background()
        add_text(slide, line, x + 0.18, top, w - 0.18, row_h, size=size, color=color)


def add_header(slide, title: str, subtitle: str, theme: Theme, page: int) -> None:
    add_text(slide, title, 0.65, 0.42, 9.7, 0.5, size=24, color=theme.ink, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.68, 0.94, 10.8, 0.3, size=9.5, color=theme.muted)
    add_text(slide, f"{page:02d}", 12.1, 0.42, 0.55, 0.32, size=10, color=theme.muted, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.68), Inches(1.28), Inches(11.97), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(theme.pale)
    line.line.fill.background()


def add_footer(slide, source: str, theme: Theme) -> None:
    add_text(slide, source, 0.68, 7.08, 11.9, 0.2, size=7.5, color=theme.muted)


def add_panel(slide, x: float, y: float, w: float, h: float, fill: str, line: str | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_metric(slide, label: str, value: str, delta: str, x: float, y: float, w: float, theme: Theme) -> None:
    add_text(slide, value, x, y, w, 0.55, size=29, color=theme.ink, bold=True)
    add_text(slide, label, x, y + 0.62, w, 0.28, size=10, color=theme.muted)
    color = theme.secondary if str(delta).startswith("+") else theme.warning if str(delta).startswith("-") else theme.muted
    add_text(slide, str(delta), x, y + 0.92, w, 0.24, size=9, color=color, bold=True)


def fit_text_size(text: str, base: float, *, long_at: int = 70, floor: float = 11) -> float:
    length = len(str(text))
    if length <= long_at:
        return base
    return max(floor, base * long_at / length)


def save(prs: Presentation, path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return path

