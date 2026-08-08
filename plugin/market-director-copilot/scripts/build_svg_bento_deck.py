from __future__ import annotations

import argparse
import json
import math
from pathlib import Path
from typing import Any
from xml.sax.saxutils import escape


WIDTH = 1280
HEIGHT = 720
MARGIN = 48
GAP = 20
FONT = "Microsoft YaHei, Arial, sans-serif"

THEMES = {
    "ceo-weekly": {"background": "#F7F7F5", "ink": "#151719", "muted": "#667079", "accent": "#C92A3A", "secondary": "#147D64", "pale": "#E8ECEA", "card": "#FFFFFF"},
    "government-formal": {"background": "#FAF9F7", "ink": "#1A1A1A", "muted": "#6D625E", "accent": "#A61B1B", "secondary": "#B38A2E", "pale": "#EFE9E5", "card": "#FFFFFF"},
    "customer-solution": {"background": "#F8FAFB", "ink": "#172126", "muted": "#65747C", "accent": "#147D64", "secondary": "#2563A6", "pale": "#E7EFEE", "card": "#FFFFFF"},
    "industry-research": {"background": "#FAFAF8", "ink": "#1A2421", "muted": "#66706C", "accent": "#167C6A", "secondary": "#D49B24", "pale": "#E7EFE9", "card": "#FFFFFF"},
    "sales-review": {"background": "#F8F9F7", "ink": "#17201C", "muted": "#67706B", "accent": "#1F8066", "secondary": "#D44A4A", "pale": "#E6ECE8", "card": "#FFFFFF"},
}


def read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def clean_color(value: Any, fallback: str) -> str:
    value = str(value or fallback)
    if not value.startswith("#"):
        value = f"#{value}"
    if len(value) != 7 or any(char not in "0123456789abcdefABCDEF" for char in value[1:]):
        return fallback
    return value


def wrap_text(value: Any, max_chars: int) -> list[str]:
    text = str(value or "").strip()
    if not text:
        return []
    lines: list[str] = []
    for paragraph in text.splitlines():
        paragraph = paragraph.strip()
        while len(paragraph) > max_chars:
            cut = paragraph.rfind(" ", 0, max_chars + 1)
            if cut < max_chars // 2:
                cut = max_chars
            lines.append(paragraph[:cut].strip())
            paragraph = paragraph[cut:].strip()
        if paragraph:
            lines.append(paragraph)
    return lines or [""]


def svg_text(value: Any, x: float, y: float, width: float, *, size: int, color: str, weight: int = 400, max_lines: int | None = None, anchor: str = "start") -> str:
    max_chars = max(8, int(width / max(1, size * 0.95)))
    lines = wrap_text(value, max_chars)
    clipped = False
    if max_lines and len(lines) > max_lines:
        lines = lines[:max_lines]
        lines[-1] = lines[-1].rstrip("，。；、 ") + "…"
        clipped = True
    if not lines:
        return ""
    clipped_attr = ' data-clipped="true"' if clipped else ""
    attrs = f'x="{x:.1f}" y="{y:.1f}" fill="{color}" font-family="{FONT}" font-size="{size}px" font-weight="{weight}" text-anchor="{anchor}" dominant-baseline="hanging"{clipped_attr}'
    tspans = []
    for index, line in enumerate(lines):
        dy = "0" if index == 0 else f"{size * 1.35:.1f}"
        tspans.append(f'<tspan x="{x:.1f}" dy="{dy}">{escape(line)}</tspan>')
    body = "".join(tspans)
    return f"<text {attrs}>{body}</text>"


def layout_boxes(count: int, layout: str, x: float, y: float, width: float, height: float) -> list[tuple[float, float, float, float]]:
    if count <= 0:
        return []
    if layout == "auto":
        layout = {1: "single", 2: "two", 3: "three"}.get(count, "hero-grid")
    if layout == "single" or count == 1:
        return [(x, y, width, height)]
    if layout == "hero-grid":
        hero_width = width * 0.62
        right_x = x + hero_width + GAP
        right_width = width - hero_width - GAP
        row_height = (height - GAP * (count - 2)) / max(1, count - 1)
        return [(x, y, hero_width, height)] + [(right_x, y + index * (row_height + GAP), right_width, row_height) for index in range(count - 1)]
    if layout == "mixed":
        lead_width = width * 0.56
        right_x = x + lead_width + GAP
        right_width = width - lead_width - GAP
        rest = count - 1
        cols = 2 if rest > 1 else 1
        rows = math.ceil(rest / cols)
        cell_width = (right_width - GAP * (cols - 1)) / cols
        cell_height = (height - GAP * (rows - 1)) / rows
        boxes = [(x, y, lead_width, height)]
        for index in range(rest):
            row, col = divmod(index, cols)
            boxes.append((right_x + col * (cell_width + GAP), y + row * (cell_height + GAP), cell_width, cell_height))
        return boxes
    cols = 2 if layout == "two" else 3 if layout == "three" else min(3, count)
    cols = min(cols, count)
    rows = math.ceil(count / cols)
    cell_width = (width - GAP * (cols - 1)) / cols
    cell_height = (height - GAP * (rows - 1)) / rows
    boxes = [(x + col * (cell_width + GAP), y + row * (cell_height + GAP), cell_width, cell_height) for row in range(rows) for col in range(cols)]
    return boxes[:count]


def render_card(card: dict[str, Any], box: tuple[float, float, float, float], theme: dict[str, str], index: int) -> str:
    x, y, width, height = box
    tone = str(card.get("tone") or "pale")
    fill = clean_color(card.get("fill"), theme["accent"] if tone == "accent" else theme["card"] if tone == "white" else theme["pale"])
    ink = clean_color(card.get("ink"), "#FFFFFF" if tone == "accent" else theme["ink"])
    muted = clean_color(card.get("muted"), "#F8E9E9" if tone == "accent" else theme["muted"])
    pad = 24
    parts = [f'<g data-card="{index}" data-x="{x:.1f}" data-y="{y:.1f}" data-width="{width:.1f}" data-height="{height:.1f}">']
    parts.append(f'<rect x="{x:.1f}" y="{y:.1f}" width="{width:.1f}" height="{height:.1f}" rx="18" fill="{fill}"/>')
    eyebrow = card.get("eyebrow")
    if eyebrow:
        parts.append(svg_text(eyebrow, x + pad, y + pad, width - pad * 2, size=12, color=muted, weight=600, max_lines=1))
    title_y = y + pad + (24 if eyebrow else 0)
    parts.append(svg_text(card.get("title") or "待填写", x + pad, title_y, width - pad * 2, size=20 if height > 220 else 16, color=ink, weight=700, max_lines=2))
    value = card.get("value")
    if value not in (None, ""):
        parts.append(svg_text(value, x + width - pad, title_y, width * 0.42, size=28, color=ink, weight=700, max_lines=1, anchor="end"))
    body_y = title_y + (70 if height > 220 else 54)
    available = max(1, int((height - (body_y - y) - 42) / 22))
    body_lines = wrap_text(card.get("body"), max(8, int((width - pad * 2) / 15)))
    if body_lines:
        max_body_lines = min(available, 5)
        parts.append(svg_text("\n".join(body_lines), x + pad, body_y, width - pad * 2, size=15 if height > 220 else 13, color=ink, max_lines=max_body_lines))
        body_y += min(len(body_lines), max_body_lines) * (15 if height > 220 else 13) * 1.35 + 12
    for bullet in (card.get("bullets") or [])[:3]:
        bullet_lines = wrap_text(bullet, max(8, int((width - pad * 2 - 18) / 13)))
        if bullet_lines:
            parts.append(svg_text("• " + bullet_lines[0], x + pad, body_y, width - pad * 2, size=13, color=ink, max_lines=1))
            body_y += 22
    parts.append("</g>")
    return "".join(parts)


def render_slide(slide: dict[str, Any], theme: dict[str, str], page: int, total: int) -> tuple[str, list[dict[str, float]]]:
    title = slide.get("title") or "待填写"
    subtitle = slide.get("subtitle") or ""
    layout = str(slide.get("layout") or "auto")
    cards = slide.get("cards") or []
    boxes = layout_boxes(len(cards), layout, MARGIN, 150, WIDTH - MARGIN * 2, 510)
    parts = [f'<svg xmlns="http://www.w3.org/2000/svg" width="{WIDTH}" height="{HEIGHT}" viewBox="0 0 {WIDTH} {HEIGHT}" role="img" aria-labelledby="title-{page}">']
    parts.append(f'<title id="title-{page}">{escape(str(title))}</title>')
    parts.append(f'<rect width="{WIDTH}" height="{HEIGHT}" fill="{theme["background"]}"/>')
    parts.append(f'<rect x="0" y="0" width="{WIDTH}" height="10" fill="{theme["accent"]}"/>')
    parts.append(svg_text(title, MARGIN, 34, WIDTH - 2 * MARGIN - 80, size=30, color=theme["ink"], weight=700, max_lines=1))
    if subtitle:
        parts.append(svg_text(subtitle, MARGIN, 82, WIDTH - 2 * MARGIN - 80, size=14, color=theme["muted"], max_lines=1))
    parts.append(f'<text x="{WIDTH - MARGIN}" y="42" fill="{theme["muted"]}" font-family="{FONT}" font-size="12px" text-anchor="end" dominant-baseline="hanging">{page:02d}/{total:02d}</text>')
    parts.append(f'<rect x="{MARGIN}" y="124" width="{WIDTH - 2 * MARGIN}" height="2" fill="{theme["pale"]}"/>')
    for index, (card, box) in enumerate(zip(cards, boxes), 1):
        parts.append(render_card(card, box, theme, index))
    source = slide.get("source") or "来源：待补充"
    parts.append(svg_text(source, MARGIN, 684, WIDTH - 2 * MARGIN, size=10, color=theme["muted"], max_lines=1))
    parts.append("</svg>")
    manifest_boxes = [{"x": round(x, 1), "y": round(y, 1), "width": round(w, 1), "height": round(h, 1)} for x, y, w, h in boxes]
    return "".join(parts), manifest_boxes


class SvgImage:
    """Small adapter so python-pptx can package an SVG with explicit dimensions."""

    def __init__(self, blob: bytes, filename: str):
        self.blob = blob
        self.filename = filename

    @property
    def ext(self) -> str:
        return "svg"

    @property
    def content_type(self) -> str:
        return "image/svg+xml"

    @classmethod
    def from_file(cls, path: Path) -> "SvgImage":
        return cls(path.read_bytes(), path.name)


def build_svg_pptx(svg_paths: list[Path], output: Path, title: str) -> Path:
    from pptx import Presentation
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.parts.image import ImagePart
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(WIDTH / 96)
    prs.slide_height = Inches(HEIGHT / 96)
    blank = prs.slide_layouts[6]
    for svg_path in svg_paths:
        slide = prs.slides.add_slide(blank)
        image_part = ImagePart.new(prs.part.package, SvgImage.from_file(svg_path))
        rel_id = slide.part.relate_to(image_part, RT.IMAGE)
        # python-pptx asks Pillow to inspect image dimensions; SVG is deliberately
        # packaged with the known 1280x720 dimensions instead.
        shape_id = slide.shapes._next_shape_id
        slide.shapes._grpSp.add_pic(shape_id, f"SVG page {shape_id - 1}", image_part.desc, rel_id, 0, 0, prs.slide_width, prs.slide_height)
    prs.core_properties.title = title
    prs.core_properties.subject = "SVG-first Bento Grid presentation"
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


def build(input_path: Path, svg_dir: Path, pptx_path: Path | None = None) -> dict[str, Any]:
    payload = read_json(input_path)
    deck = payload.get("deck") or {}
    width = int(deck.get("width", WIDTH))
    height = int(deck.get("height", HEIGHT))
    if (width, height) != (WIDTH, HEIGHT):
        raise ValueError("SVG-first canvas must be exactly 1280x720")
    theme = THEMES.get(str(deck.get("theme") or "industry-research"), THEMES["industry-research"])
    slides = payload.get("slides") or []
    if not slides:
        raise ValueError("slides must contain at least one page")
    svg_dir.mkdir(parents=True, exist_ok=True)
    svg_paths: list[Path] = []
    manifest_pages = []
    for page, slide in enumerate(slides, 1):
        svg, boxes = render_slide(slide, theme, page, len(slides))
        path = svg_dir / f"slide-{page:02d}.svg"
        path.write_text(svg, encoding="utf-8")
        svg_paths.append(path)
        manifest_pages.append({"page": page, "title": slide.get("title", ""), "layout": slide.get("layout", "auto"), "boxes": boxes, "file": path.name})
    manifest = {"mode": "svg-first", "canvas": {"width": WIDTH, "height": HEIGHT}, "theme": deck.get("theme", "industry-research"), "pages": manifest_pages}
    (svg_dir / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    if pptx_path:
        build_svg_pptx(svg_paths, pptx_path, str(deck.get("title") or "SVG-first Bento Grid"))
    return {"svg_dir": str(svg_dir), "pages": len(svg_paths), "pptx": str(pptx_path) if pptx_path else None}


def main() -> int:
    parser = argparse.ArgumentParser(description="Build SVG-first Bento Grid presentation pages.")
    parser.add_argument("--input", type=Path, required=True)
    parser.add_argument("--svg-dir", type=Path, required=True)
    parser.add_argument("--pptx", type=Path)
    args = parser.parse_args()
    print(json.dumps(build(args.input.resolve(), args.svg_dir.resolve(), args.pptx.resolve() if args.pptx else None), ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
