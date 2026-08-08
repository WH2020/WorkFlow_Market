from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from deck_utils import add_footer, add_header, add_lines, add_metric, add_panel, add_text, new_presentation, rgb, save, set_background, THEMES


def text(value, fallback="待填写"):
    if value in (None, "", [], {}):
        return fallback
    return str(value)


def build(input_path: Path, output_path: Path, template_id: str = "ceo-weekly") -> Path:
    payload = json.loads(input_path.read_text(encoding="utf-8"))
    theme = THEMES.get(template_id, THEMES["ceo-weekly"])
    prs = new_presentation()
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_background(slide, theme.background)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, 160000)
    band.fill.solid(); band.fill.fore_color.rgb = rgb(theme.accent); band.line.fill.background()
    add_text(slide, "市场工作周报", 0.75, 1.28, 8.0, 0.7, size=34, color=theme.ink, bold=True)
    add_text(slide, text(payload.get("reporting_period")), 0.78, 2.15, 7.5, 0.45, size=16, color=theme.muted)
    add_text(slide, text(payload.get("headline")), 0.78, 3.0, 10.8, 0.72, size=25, color=theme.ink, bold=True)
    add_text(slide, "汇报对象：CEO", 0.8, 5.85, 4.0, 0.3, size=10, color=theme.muted)
    add_text(slide, "MARKET DIRECTOR", 10.0, 5.85, 2.6, 0.3, size=9, color=theme.accent, bold=True, align=PP_ALIGN.RIGHT)

    # 2. Executive summary and decisions
    slide = prs.slides.add_slide(blank); set_background(slide, theme.background)
    add_header(slide, "本周核心结论", "CEO 需要关注的判断和决策", theme, 2)
    add_panel(slide, 0.68, 1.65, 7.65, 4.95, "FFFFFF", theme.pale)
    add_text(slide, "结论", 1.0, 1.98, 6.4, 0.35, size=16, color=theme.ink, bold=True)
    add_lines(slide, payload.get("executive_summary", []), 1.0, 2.58, 6.7, 2.7, size=14, color=theme.ink, bullet_color=theme.accent)
    add_panel(slide, 8.6, 1.65, 4.05, 4.95, theme.pale)
    add_text(slide, "需要 CEO 决策", 8.95, 1.98, 3.2, 0.35, size=15, color=theme.ink, bold=True)
    decisions = [f"{item.get('title', '事项')}：{item.get('detail', '')}" for item in payload.get("decision_requests", [])]
    add_lines(slide, decisions, 8.95, 2.58, 3.25, 3.1, size=11.5, color=theme.ink, bullet_color=theme.warning)
    add_footer(slide, "来源：本周工作台记录；待确认事项已标注", theme)

    # 3. Metrics
    slide = prs.slides.add_slide(blank); set_background(slide, theme.background)
    add_header(slide, "本周成果与关键指标", "将工作量转化为对业务有意义的进展", theme, 3)
    metrics = payload.get("metrics", [])[:4]
    for idx, item in enumerate(metrics):
        x = 0.75 + idx * 3.05
        add_panel(slide, x, 1.75, 2.7, 2.05, theme.pale)
        add_metric(slide, text(item.get("label")), text(item.get("value")), text(item.get("delta"), "0"), x + 0.28, 2.05, 2.15, theme)
    add_panel(slide, 0.75, 4.25, 11.85, 2.0, "FFFFFF", theme.pale)
    add_text(slide, "本周工作结果", 1.05, 4.58, 3.4, 0.3, size=15, color=theme.ink, bold=True)
    add_lines(slide, payload.get("executive_summary", []), 1.05, 5.05, 10.9, 0.9, size=12.5, color=theme.ink, bullet_color=theme.secondary)
    add_footer(slide, "来源：周报输入与销售台账", theme)

    # 4. Intelligence
    slide = prs.slides.add_slide(blank); set_background(slide, theme.background)
    add_header(slide, "行业、竞争与政策情报", "脑机、具身、数采及相邻方向的商业含义", theme, 4)
    findings = payload.get("intelligence", [])[:3]
    for idx, item in enumerate(findings):
        x = 0.75 + idx * 4.05
        add_panel(slide, x, 1.7, 3.65, 4.75, theme.pale)
        add_text(slide, text(item.get("topic")), x + 0.3, 2.05, 2.8, 0.42, size=19, color=theme.accent, bold=True)
        add_text(slide, text(item.get("finding")), x + 0.3, 2.78, 3.0, 1.35, size=13, color=theme.ink)
        add_text(slide, "商业含义", x + 0.3, 4.45, 2.8, 0.3, size=11, color=theme.muted, bold=True)
        add_text(slide, text(item.get("implication")), x + 0.3, 4.83, 3.0, 0.75, size=15, color=theme.ink, bold=True)
        add_text(slide, f"来源：{text(item.get('source'))}", x + 0.3, 5.8, 3.0, 0.3, size=8.5, color=theme.muted)
    add_footer(slide, "关键结论必须回溯到来源登记表", theme)

    # 5. Pipeline
    slide = prs.slides.add_slide(blank); set_background(slide, theme.background)
    add_header(slide, "销售漏斗和阶段变化", "只使用已有客户证据，不用主观热度代替进展", theme, 5)
    stages = payload.get("pipeline", [])
    max_count = max([int(item.get("count", 0)) for item in stages] or [1])
    for idx, item in enumerate(stages[:6]):
        y = 1.75 + idx * 0.72
        label = text(item.get("stage"))
        count = int(item.get("count", 0))
        add_text(slide, label, 0.9, y, 2.3, 0.28, size=12, color=theme.ink)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 3.1 * 914400, int(y * 914400), int(max(0.3, 7.8 * count / max_count) * 914400), int(0.32 * 914400))
        bar.fill.solid(); bar.fill.fore_color.rgb = rgb(theme.secondary if idx < 3 else theme.accent); bar.line.fill.background()
        add_text(slide, str(count), 11.25, y - 0.01, 0.7, 0.3, size=13, color=theme.ink, bold=True, align=PP_ALIGN.RIGHT)
    add_panel(slide, 0.85, 6.1, 11.7, 0.45, theme.pale)
    add_text(slide, "管理提示：阶段变化必须关联客户、证据、下一步和截止日期。", 1.05, 6.18, 11.1, 0.25, size=10.5, color=theme.ink)
    add_footer(slide, "来源：data/sales/customers.csv 与 activities.csv", theme)

    # 6. Customers
    slide = prs.slides.add_slide(blank); set_background(slide, theme.background)
    add_header(slide, "重点客户进展、风险与下一步", "红黄绿是管理提示，不是对客户或销售的结论", theme, 6)
    customers = payload.get("priority_customers", [])[:3]
    for idx, item in enumerate(customers):
        x = 0.75 + idx * 4.05
        health = text(item.get("health"), "黄")
        health_color = theme.secondary if health == "绿" else theme.warning if health == "黄" else theme.accent
        add_panel(slide, x, 1.7, 3.65, 4.85, "FFFFFF", theme.pale)
        add_text(slide, text(item.get("name")), x + 0.3, 2.05, 2.65, 0.38, size=16, color=theme.ink, bold=True)
        add_text(slide, health, x + 3.0, 2.03, 0.35, 0.35, size=13, color=health_color, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, f"阶段：{text(item.get('stage'))}", x + 0.3, 2.78, 3.0, 0.3, size=11, color=theme.muted)
        add_text(slide, "进展", x + 0.3, 3.45, 3.0, 0.25, size=10, color=theme.muted, bold=True)
        add_text(slide, text(item.get("progress")), x + 0.3, 3.75, 3.0, 0.65, size=12.5, color=theme.ink)
        add_text(slide, "风险", x + 0.3, 4.65, 3.0, 0.25, size=10, color=theme.muted, bold=True)
        add_text(slide, text(item.get("risk")), x + 0.3, 4.95, 3.0, 0.55, size=12.5, color=theme.ink)
        add_text(slide, "下一步", x + 0.3, 5.68, 3.0, 0.25, size=10, color=theme.muted, bold=True)
        add_text(slide, text(item.get("next_action")), x + 0.3, 5.98, 3.0, 0.4, size=11.5, color=theme.ink, bold=True)
    add_footer(slide, "来源：客户活动证据与销售沟通复盘", theme)

    # 7. Sales and government
    slide = prs.slides.add_slide(blank); set_background(slide, theme.background)
    add_header(slide, "销售复盘、政府合作与资源问题", "把发现转成可执行的管理动作", theme, 7)
    add_panel(slide, 0.75, 1.7, 5.8, 4.95, theme.pale)
    add_text(slide, "销售复盘", 1.05, 2.02, 3.2, 0.35, size=16, color=theme.ink, bold=True)
    review_lines = [f"{item.get('salesperson', '销售')}：{item.get('gap', '')}；管理动作：{item.get('manager_action', '')}" for item in payload.get("sales_review", [])]
    add_lines(slide, review_lines, 1.05, 2.68, 5.0, 3.2, size=11.5, color=theme.ink, bullet_color=theme.accent)
    add_panel(slide, 6.85, 1.7, 5.75, 4.95, "FFFFFF", theme.pale)
    add_text(slide, "政府合作与资源", 7.15, 2.02, 3.8, 0.35, size=16, color=theme.ink, bold=True)
    gov_lines = [f"{item.get('region', '地区')}：{item.get('scenario', '')}，状态 {item.get('status', '')}，下一步 {item.get('next_action', '')}" for item in payload.get("government_projects", [])]
    add_lines(slide, gov_lines, 7.15, 2.68, 4.95, 3.2, size=11.5, color=theme.ink, bullet_color=theme.secondary)
    add_footer(slide, "资源申请需补充业务理由、截止时间和预期结果", theme)

    # 8. Next week
    slide = prs.slides.add_slide(blank); set_background(slide, theme.background)
    add_header(slide, "下周重点", "责任人、截止时间和成功标准", theme, 8)
    add_panel(slide, 0.75, 1.7, 8.0, 4.95, "FFFFFF", theme.pale)
    add_lines(slide, payload.get("next_week", []), 1.1, 2.25, 7.1, 2.8, size=16, color=theme.ink, bullet_color=theme.accent)
    add_panel(slide, 9.0, 1.7, 3.6, 4.95, theme.pale)
    add_text(slide, "复盘机制", 9.3, 2.1, 2.7, 0.35, size=15, color=theme.ink, bold=True)
    add_lines(slide, ["周一：确认优先级", "周三：检查承诺", "周五：沉淀证据"], 9.3, 2.8, 2.85, 2.5, size=12, color=theme.ink, bullet_color=theme.secondary)
    add_footer(slide, "来源：周报输入、知识库、客户台账和已确认沟通记录", theme)

    prs.core_properties.title = f"市场工作周报 - {text(payload.get('reporting_period'))}"
    prs.core_properties.subject = "CEO 市场周报"
    prs.core_properties.author = "市场总监工作台"
    return save(prs, output_path)


def main() -> int:
    parser = argparse.ArgumentParser(description="Build an editable 8-slide CEO weekly report.")
    parser.add_argument("--input", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--template", default="ceo-weekly")
    args = parser.parse_args()
    result = build(args.input.resolve(), args.output.resolve(), args.template)
    print(result)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

